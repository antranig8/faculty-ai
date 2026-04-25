from io import BytesIO

from app.models.response_models import Slide
from app.services.presentation_preparer import _build_slide


def parse_pptx_slides(file_bytes: bytes) -> list[Slide]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX parsing requires python-pptx. Install backend requirements first.") from exc

    presentation = Presentation(BytesIO(file_bytes))
    slides: list[Slide] = []

    for index, pptx_slide in enumerate(presentation.slides, start=1):
        text_runs: list[str] = []
        for shape in pptx_slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_runs.append(text)

        title = text_runs[0].splitlines()[0] if text_runs else f"Slide {index}"
        content = "\n".join(text_runs[1:] if len(text_runs) > 1 else text_runs)
        slides.append(_build_slide(index, title, content))

    return slides

